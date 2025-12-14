"""Document processing service for extracting, chunking, and summarizing documents."""

import asyncio
from pathlib import Path
from typing import List, Optional, Tuple
from datetime import datetime
from uuid import UUID, uuid4
import tiktoken

from pypdf import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation

from backend.models.document import DocumentChunk, ProcessedDocument, FileType, ChunkingStrategy, DocumentMetadata
from backend.config import settings


class DocumentProcessor:
    def __init__(self, embedding_service, chat_service):
        self.embedding_service = embedding_service
        self.chat_service = chat_service
        self.tokenizer = tiktoken.get_encoding("cl100k_base")

    async def process_document(self, file_path: Path) -> Optional[ProcessedDocument]:
        """Process a single document: extract, summarize, chunk."""
        try:
            file_path = Path(file_path)
            if not file_path.exists():
                print(f"❌ File not found: {file_path}")
                return None

            # Extract based on file type
            if file_path.suffix.lower() == ".pdf":
                full_content, page_texts, author, created_at = self.extract_pdf(file_path)
                file_type = FileType.PDF
            elif file_path.suffix.lower() == ".docx":
                full_content, page_texts, author, created_at = self.extract_docx(file_path)
                file_type = FileType.DOCX
            elif file_path.suffix.lower() == ".pptx":
                full_content, page_texts, author, created_at = self.extract_pptx(file_path)
                file_type = FileType.PPTX
            else:
                print(f"❌ Unsupported file type: {file_path.suffix}")
                return None

            # Generate summary and keywords from complete content
            summary, keywords = await self.generate_summary_and_keywords(full_content, file_path.stem)

            # Create document ID
            document_id = uuid4()

            # Create metadata
            metadata = DocumentMetadata(
                title=file_path.stem,
                author=author,
                created_at=created_at,
                file_type=file_type,
                file_path=str(file_path),
                file_size=file_path.stat().st_size,
            )

            # Chunk using all three strategies
            whole_file_chunks = self.chunk_whole_file(document_id, full_content, metadata, summary, keywords)
            page_chunks = self.chunk_by_pages(document_id, page_texts, metadata, summary, keywords)
            token_chunks = self.chunk_by_max_tokens(document_id, full_content, metadata, summary, keywords)

            return ProcessedDocument(
                document_id=document_id,
                metadata=metadata,
                full_content=full_content,
                whole_file_chunks=whole_file_chunks,
                page_chunks=page_chunks,
                token_chunks=token_chunks,
                summary=summary,
                keywords=keywords,
            )

        except Exception as e:
            print(f"❌ Error processing {file_path.name}: {e}")
            return None

    def extract_pdf(self, file_path: Path) -> Tuple[str, List[str], Optional[str], Optional[datetime]]:
        """Extract text and metadata from PDF file."""
        reader = PdfReader(file_path)

        # Extract text from all pages
        page_texts = [page.extract_text() or "" for page in reader.pages]
        full_text = "\n\n".join(page_texts)

        # Extract metadata
        metadata = reader.metadata
        author = metadata.get("/Author") if metadata else None
        created_at = None

        if metadata and "/CreationDate" in metadata:
            date_str = metadata["/CreationDate"]
            if date_str.startswith("D:"):
                try:
                    created_at = datetime.strptime(date_str[2:16], "%Y%m%d%H%M%S")
                except ValueError:
                    created_at = None

        return full_text, page_texts, author, created_at

    def extract_docx(self, file_path: Path) -> Tuple[str, List[str], Optional[str], Optional[datetime]]:
        """Extract text and metadata from DOCX file."""
        doc = DocxDocument(file_path)

        # Extract text from all paragraphs
        full_text = "\n".join([para.text for para in doc.paragraphs])
        page_texts = [full_text]  # DOCX doesn't have page concept, so single "page"

        # Extract metadata
        core_properties = doc.core_properties
        author = core_properties.author
        created_at = core_properties.created

        return full_text, page_texts, author, created_at

    def extract_pptx(self, file_path: Path) -> Tuple[str, List[str], Optional[str], Optional[datetime]]:
        """Extract text and metadata from PPTX file."""
        prs = Presentation(file_path)

        # Extract text from all slides
        slide_texts = []
        for slide in prs.slides:
            slide_text = "\n".join([
                shape.text for shape in slide.shapes 
                if hasattr(shape, "text") and shape.text.strip()
            ])
            slide_texts.append(slide_text)

        full_text = "\n\n".join(slide_texts)

        # Extract metadata
        core_properties = prs.core_properties
        author = core_properties.author
        created_at = core_properties.created

        return full_text, slide_texts, author, created_at

    async def generate_summary_and_keywords(self, full_content: str, title: str) -> Tuple[str, List[str]]:
        """Generate summary and keywords using separate functions."""
        summary = await self.generate_summary(full_content, title)
        keywords = await self.extract_keywords(full_content, title)
        return summary, keywords

    async def generate_summary(self, full_content: str, title: str) -> str:
        """Generate a concise summary of the document."""
        prompt = f"""Generate a concise 2-3 sentence summary for this research document titled '{title}'.

Document content (first 8000 characters):
{full_content[:8000]}

Instructions:
- Focus on the main research purpose, key findings, or conclusions
- Avoid just repeating the title or listing author names
- Keep it professional, accurate, and substantive
- Write in complete sentences

Provide only the summary, no labels or prefixes."""

        try:
            response = await self.chat_service.generate_completion(prompt, max_tokens=150, temperature=0.3)
            return response.strip()
        except Exception as e:
            print(f"Error generating summary: {e}")
            return f"Research document analyzing {title.replace('_', ' ').lower()} with clinical findings and analysis."

    async def extract_keywords(self, full_content: str, title: str) -> List[str]:
        """Extract relevant keywords from the document."""
        prompt = f"""Extract 5-7 key terms/keywords from this research document titled '{title}'.

Document content (first 4000 characters):
{full_content[:4000]}

Instructions:
- Focus on technical terms, medical concepts, methodologies, and domain-specific terminology
- Include relevant acronyms (e.g., RCT, MRI) but NOT reference identifiers (PMID, DOI)
- Avoid generic terms like "study", "research", "data" unless they're part of a specific concept
- Each keyword should be 1-4 words maximum
- DO NOT include author names or citation elements

Provide ONLY the keywords as a comma-separated list, nothing else."""

        try:
            response = await self.chat_service.generate_completion(prompt, max_tokens=60, temperature=0.2)
            # Parse comma-separated keywords
            keywords = [k.strip() for k in response.split(",") if k.strip()]
            # Clean up keywords: remove empty strings, normalize whitespace
            keywords = [' '.join(k.split()) for k in keywords if k.strip()]
            return keywords[:7] if len(keywords) >= 7 else keywords
        except Exception as e:
            print(f"Error extracting keywords: {e}")
            return ["clinical research", "medical study", "health outcomes", "data analysis"]

    def chunk_whole_file(self, document_id: UUID, content: str, metadata: DocumentMetadata,
                        summary: str, keywords: List[str]) -> List[DocumentChunk]:
        """Chunk entire document as single chunk."""
        token_count = len(self.tokenizer.encode(content))

        return [DocumentChunk(
            document_id=document_id,
            content=content,
            chunk_index=0,
            chunking_strategy=ChunkingStrategy.WHOLE_FILE,
            token_count=token_count,
            title=metadata.title,
            author=metadata.author,
            created_at=metadata.created_at,
            file_type=metadata.file_type,
            file_path=metadata.file_path,
            summary=summary,
            keywords=keywords,
        )]

    def chunk_by_pages(self, document_id: UUID, page_texts: List[str], metadata: DocumentMetadata,
                      summary: str, keywords: List[str]) -> List[DocumentChunk]:
        """Chunk document by pages/slides."""
        chunks = []

        for idx, page_text in enumerate(page_texts):
            if not page_text.strip():
                continue

            token_count = len(self.tokenizer.encode(page_text))

            chunks.append(DocumentChunk(
                document_id=document_id,
                content=page_text,
                chunk_index=idx,
                chunking_strategy=ChunkingStrategy.PAGES,
                token_count=token_count,
                title=metadata.title,
                author=metadata.author,
                created_at=metadata.created_at,
                file_type=metadata.file_type,
                file_path=metadata.file_path,
                summary=summary,
                keywords=keywords,
            ))

        return chunks

    def chunk_by_max_tokens(self, document_id: UUID, content: str, metadata: DocumentMetadata,
                           summary: str, keywords: List[str], max_tokens: int = 512) -> List[DocumentChunk]:
        """Chunk document by maximum token size with overlap."""
        overlap_tokens = max_tokens // 4  # 25% overlap
        tokens = self.tokenizer.encode(content)
        chunks = []
        idx = 0

        while idx < len(tokens):
            chunk_tokens = tokens[idx : idx + max_tokens]
            chunk_text = self.tokenizer.decode(chunk_tokens)

            chunks.append(DocumentChunk(
                document_id=document_id,
                content=chunk_text,
                chunk_index=len(chunks),
                chunking_strategy=ChunkingStrategy.MAX_TOKENS,
                token_count=len(chunk_tokens),
                title=metadata.title,
                author=metadata.author,
                created_at=metadata.created_at,
                file_type=metadata.file_type,
                file_path=metadata.file_path,
                summary=summary,
                keywords=keywords,
            ))

            idx += max_tokens - overlap_tokens

        return chunks