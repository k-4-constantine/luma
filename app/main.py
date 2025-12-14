import os
import re
import fitz # PyMuPDF
import glob
import json
import requests
import numpy as np
import concurrent.futures
import docx  # NEW: Import python-docx
from datetime import datetime
from typing import List, Dict, Any
from fastapi import FastAPI, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.staticfiles import StaticFiles
from pydantic import BaseModel
from pptx import Presentation
import chromadb
from chromadb.api.types import Documents, EmbeddingFunction, Embeddings
from dotenv import load_dotenv
from sklearn.metrics.pairwise import cosine_similarity

# --- Configuration ---
load_dotenv()
API_KEY = os.getenv("GREENPT_API_KEY")
API_BASE = os.getenv("GREENPT_API_URL", "https://api.greenpt.ai/v1").rstrip('/')
CHAT_MODEL = os.getenv("GREENPT_MODEL", "green-l-raw") 

# Define chunk size (in characters) for processing Word documents as streaming documents
DOC_CHUNK_SIZE = 800  
DOC_CHUNK_OVERLAP = 100 # Simple overlap to prevent context cutoff

if not API_KEY:
    print("FATAL: GREENPT_API_KEY not set.")

# --- GreenPT Embedding Function ---
class GreenPTEmbeddingFunction(EmbeddingFunction):
    def __call__(self, input: Documents) -> Embeddings:
        url = f"{API_BASE}/embeddings"
        payload = {
            "model": "green-embedding",
            "input": input,
            "encoding_format": "float"
        }
        headers = {"Authorization": f"Bearer {API_KEY}", "Content-Type": "application/json"}
        
        for attempt in range(3):
            try:
                resp = requests.post(url, json=payload, headers=headers, timeout=30)
                if resp.status_code == 200:
                    data = resp.json()
                    sorted_data = sorted(data['data'], key=lambda x: x['index'])
                    return [item['embedding'] for item in sorted_data]
                elif resp.status_code == 429:
                    import time
                    time.sleep(2 * (attempt + 1))
                    continue
            except Exception:
                import time
                time.sleep(1)
        raise Exception("Failed to generate embeddings after retries")

# --- Database Init ---
client = chromadb.PersistentClient(path="./chroma_data")
try:
    collection = client.get_collection(name="simple_rag")
except:
    collection = client.create_collection(name="simple_rag", embedding_function=GreenPTEmbeddingFunction())

app = FastAPI()

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

if not os.path.exists("static"):
    os.makedirs("static")
app.mount("/static", StaticFiles(directory="static"), name="static")

class AskReq(BaseModel):
    question: str

# --- Helper: Generate Metadata ---
def generate_file_metadata(full_text: str) -> Dict[str, str]:
    truncated_text = full_text[:4000]
    system_prompt = (
        "You are a helpful research assistant. "
        "Analyze the provided document text. "
        "Output a JSON object with two keys: "
        "'summary' (concise summary, max 50 words) and "
        "'keywords' (comma-separated string of top 5 keywords). "
        "Do not include markdown formatting."
    )
    
    payload = {
        "model": CHAT_MODEL,
        "messages": [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": truncated_text}
        ],
        "stream": False,
        "response_format": {"type": "json_object"}
    }
    headers = {"Authorization": f"Bearer {API_KEY}", "Content-Type": "application/json"}

    try:
        resp = requests.post(f"{API_BASE}/chat/completions", json=payload, headers=headers, timeout=40)
        resp.raise_for_status()
        content = resp.json()['choices'][0]['message']['content']
        data = json.loads(content)
        return {
            "summary": data.get("summary", "N/A"),
            "keywords": data.get("keywords", "N/A")
        }
    except Exception as e:
        print(f"Metadata Gen Error: {e}")
        return {"summary": "Error generating summary", "keywords": "Error"}

# --- Helper: Text Cleaning ---
def clean_text(text: str) -> str:
    if not text: return ""
    return re.sub(r'\s+', ' ', text).strip()

# --- Helper: File Processors ---
def process_pdf_file(fpath):
    try:
        doc = fitz.open(fpath)
        meta = doc.metadata
        author = meta.get('author', 'Unknown Author') if meta else "Unknown Author"
        created_at = meta.get('creationDate', 'Unknown Date') if meta else "Unknown Date"
        
        all_text_list = []
        slides_content = []
        for i, page in enumerate(doc):
            cleaned = clean_text(page.get_text())
            if len(cleaned) > 10:
                all_text_list.append(cleaned)
                slides_content.append({"text": cleaned, "page": i + 1, "source_type": "PDF"})
        
        return {
            "author": author, "created_at": created_at,
            "full_text": "\n\n".join(all_text_list), "chunks": slides_content
        }
    except Exception as e:
        print(f"Error PDF {fpath}: {e}")
        return None

def process_pptx_file(fpath):
    try:
        prs = Presentation(fpath)
        core = prs.core_properties
        author = core.author if core.author else "Unknown Author"
        created_at = core.created.strftime("%Y-%m-%d") if core.created else "Unknown Date"
        
        all_text_list = []
        slides_content = []
        for i, slide in enumerate(prs.slides):
            texts = [s.text for s in slide.shapes if hasattr(s, "text") and s.text]
            page_text = clean_text("\n".join(texts))
            if len(page_text) > 10:
                all_text_list.append(page_text)
                slides_content.append({"text": page_text, "page": i + 1, "source_type": "PPTX"})
                
        return {
            "author": author, "created_at": created_at,
            "full_text": "\n\n".join(all_text_list), "chunks": slides_content
        }
    except Exception as e:
        print(f"Error PPTX {fpath}: {e}")
        return None

# --- NEW: Process DOCX Files with Chunking ---
def process_docx_file(fpath):
    try:
        doc = docx.Document(fpath)
        core_props = doc.core_properties
        author = core_props.author if core_props.author else "Unknown Author"
        created_at = core_props.created.strftime("%Y-%m-%d") if core_props.created else "Unknown Date"
        
        all_text_list = []
        chunks = []
        
        # Temporary variable to accumulate current chunk
        current_chunk_text = ""
        chunk_counter = 1
        
        # Iterate through paragraphs for aggregation
        for para in doc.paragraphs:
            text = clean_text(para.text)
            if not text: continue
            
            all_text_list.append(text)
            
            # If adding this paragraph doesn't exceed limit, continue adding
            if len(current_chunk_text) + len(text) < DOC_CHUNK_SIZE:
                current_chunk_text += " " + text
            else:
                # Exceeded limit, package current chunk
                if current_chunk_text:
                    chunks.append({
                        "text": current_chunk_text.strip(),
                        "page": chunk_counter, # For Word, Page is actually Chunk ID
                        "source_type": "DOCX"
                    })
                    chunk_counter += 1
                
                # Start new chunk with some overlap
                # Simple approach: start from current paragraph, can be optimized later to take latter half of previous paragraph
                current_chunk_text = text

        # Process the last unpacked chunk
        if current_chunk_text:
            chunks.append({
                "text": current_chunk_text.strip(),
                "page": chunk_counter,
                "source_type": "DOCX"
            })

        return {
            "author": author,
            "created_at": created_at,
            "full_text": "\n\n".join(all_text_list),
            "chunks": chunks
        }
    except Exception as e:
        print(f"Error DOCX {fpath}: {e}")
        return None

# --- Pipeline Processor ---
def process_single_file_pipeline(fpath):
    fname = os.path.basename(fpath)
    print(f"START Processing: {fname}")
    
    extracted = None
    # Route by file extension
    if fpath.lower().endswith(".pdf"):
        extracted = process_pdf_file(fpath)
    elif fpath.lower().endswith((".pptx", ".ppt")):
        extracted = process_pptx_file(fpath)
    elif fpath.lower().endswith(".docx"): # Add docx support
        extracted = process_docx_file(fpath)
    
    # Ignore unsupported files or files that failed extraction
    if not extracted or not extracted.get('chunks'):
        return 0
    
    ai_meta = generate_file_metadata(extracted['full_text'])
    
    docs = []
    metas = []
    ids = []
    
    for item in extracted['chunks']:
        meta = {
            "source": fname,
            "page": item['page'],
            "author": extracted['author'],
            "created_at": extracted['created_at'],
            "file_type": item['source_type'],
            "summary": ai_meta['summary'],
            "keywords": ai_meta['keywords']
        }
        docs.append(item['text'])
        metas.append(meta)
        ids.append(f"{fname}_{item['page']}")
    
    if docs:
        try:
            collection.add(documents=docs, metadatas=metas, ids=ids)
            print(f"DONE Processing: {fname} ({len(docs)} chunks)")
            return len(docs)
        except Exception as e:
            print(f"DB Error {fname}: {e}")
            return 0
    return 0

# --- Ingest Endpoint ---
@app.post("/ingest")
def ingest():
    data_dir = "./data"
    if os.path.exists("/app_data"): data_dir = "/app_data"
    elif os.path.exists("data"): data_dir = "data"
    
    # Scan all supported formats
    extensions = ["*.pptx", "*.pdf", "*.docx"]
    all_files = []
    for ext in extensions:
        all_files.extend(glob.glob(f"{data_dir}/{ext}"))

    if not all_files:
        return {"error": f"No files found in {data_dir}"}

    print(f"Found {len(all_files)} files. Starting parallel ingestion...")
    
    with concurrent.futures.ThreadPoolExecutor(max_workers=5) as executor:
        results = list(executor.map(process_single_file_pipeline, all_files))
    
    total_chunks = sum(results)
    processed_files = sum(1 for r in results if r > 0)

    return {
        "status": "ok", 
        "files": processed_files, 
        "chunks": total_chunks,
        "message": "Parallel ingestion complete (PDF, PPTX, DOCX)."
    }

@app.post("/ask")
def ask(req: AskReq):
    results = collection.query(query_texts=[req.question], n_results=5)
    if not results['documents'] or not results['documents'][0]:
        return {"answer": "No info found."}
    
    docs = results['documents'][0]
    metas = results['metadatas'][0]
    
    context = ""
    for i, (d, m) in enumerate(zip(docs, metas)):
        # Display Page/Chunk ID
        source_label = f"Page {m['page']}" if m['file_type'] != "DOCX" else f"Part {m['page']}"
        context += f"[Source {i+1}: {m['source']} ({source_label})]\n{d}\n\n"

    try:
        payload = {
            "model": CHAT_MODEL,
            "messages": [
                {"role": "system", "content": "Answer using the context."},
                {"role": "user", "content": f"Context:\n{context}\n\nQ: {req.question}"}
            ]
        }
        resp = requests.post(f"{API_BASE}/chat/completions", json=payload, headers={"Authorization": f"Bearer {API_KEY}"})
        return {"answer": resp.json()['choices'][0]['message']['content'], "citations": metas}
    except Exception as e:
        return {"error": str(e)}

@app.get("/graph")
def get_knowledge_graph():
    try:
        data = collection.get(include=['metadatas', 'embeddings'])
        if not data or not data.get('metadatas'):
            return {"nodes": [], "links": [], "categories": []}

        embeddings = data.get('embeddings', [])
        has_embeddings = (embeddings is not None and len(embeddings) > 0)

        files_map = {}
        embedding_dim = 1536 
        if has_embeddings:
            for emb in embeddings:
                if emb is not None and len(emb) > 0:
                    embedding_dim = len(emb)
                    break

        for i, meta in enumerate(data['metadatas']):
            fname = meta['source']
            if fname not in files_map:
                files_map[fname] = {
                    "author": meta.get('author', 'Unknown'),
                    "created_at": meta.get('created_at', 'Unknown'),
                    "keywords": set(),
                    "embeddings": [],
                    "chunk_count": 0
                }
            
            if 'keywords' in meta and meta['keywords']:
                raw_phrases = [k.strip().lower() for k in meta['keywords'].split(',') if k.strip()]
                for phrase in raw_phrases:
                    words = re.findall(r'\w+', phrase)
                    valid_words = [w for w in words if len(w) > 2] 
                    files_map[fname]['keywords'].update(valid_words)
                
            if has_embeddings and i < len(embeddings):
                curr_emb = embeddings[i]
                if curr_emb is not None and len(curr_emb) > 0:
                    files_map[fname]['embeddings'].append(curr_emb)
                
            files_map[fname]['chunk_count'] += 1

        nodes = []
        categories = []
        author_map = {} 
        file_names = list(files_map.keys())
        file_indices = {name: i for i, name in enumerate(file_names)}
        file_avg_embeddings = []
        
        current_time = datetime.now()

        for fname in file_names:
            info = files_map[fname]
            raw_author = info['author'].strip()
            
            is_unknown = raw_author.lower() in ["unknown author", "unknown", "n/a", "", "none"]
            author_key = "Unknown" if is_unknown else raw_author
            
            if author_key not in author_map:
                author_map[author_key] = len(categories)
                categories.append({"name": author_key})
            
            opacity = 1.0 
            try:
                created_dt = None
                date_str = str(info['created_at']).strip()
                for fmt in ["%Y-%m-%d", "%Y-%m-%d %H:%M:%S", "%Y/%m/%d", "%Y%m%d"]:
                    try:
                        created_dt = datetime.strptime(date_str, fmt)
                        break
                    except: pass
                
                if created_dt:
                    days_diff = (current_time - created_dt).days
                    if days_diff < 0: days_diff = 0
                    if days_diff <= 90: opacity = 1.0
                    elif days_diff <= 365: opacity = 0.8
                    else: opacity = max(0.2, 0.8 - ((days_diff - 365) / (365 * 2)))
                else:
                    opacity = 0.5 
            except:
                opacity = 0.5 

            avg_emb = np.zeros(embedding_dim)
            valid_embs = [e for e in info['embeddings'] if e is not None and len(e) == embedding_dim]
            if valid_embs:
                avg_emb = np.mean(valid_embs, axis=0)
            file_avg_embeddings.append(avg_emb)
            
            item_style = {
                "opacity": opacity,
                "borderColor": "#fff" if opacity > 0.8 else "transparent", 
                "borderWidth": 1
            }
            label_style = {"show": True, "color": "#333" if opacity > 0.7 else "#aaa"}

            if is_unknown:
                item_style["color"] = "#cccccc"
                label_style["color"] = "#999999"

            nodes.append({
                "id": str(file_indices[fname]),
                "name": fname,
                "category": author_map[author_key], 
                "value": info['chunk_count'], 
                "itemStyle": item_style,
                "symbolSize": 15, 
                "label": label_style
            })

        links = []
        connection_counts = {i: 0 for i in range(len(nodes))}
        
        sim_matrix = np.zeros((len(nodes), len(nodes)))
        if len(file_avg_embeddings) > 0:
            try:
                emb_matrix = np.vstack(file_avg_embeddings)
                sim_matrix = cosine_similarity(emb_matrix)
            except Exception as e:
                print(f"Sim Matrix Error: {e}")

        for i in range(len(file_names)):
            for j in range(i + 1, len(file_names)):
                source_file = file_names[i]
                target_file = file_names[j]
                
                kw_a = files_map[source_file]['keywords']
                kw_b = files_map[target_file]['keywords']
                
                intersection = kw_a.intersection(kw_b)
                match_count = len(intersection)
                
                raw_score = sim_matrix[i][j]
                sem_score = float(raw_score) 
                
                if match_count > 0:
                    width = 1 + (match_count * 0.5) 
                    links.append({
                        "source": str(i), "target": str(j),
                        "value": match_count,
                        "lineStyle": {
                            "width": min(width, 8), 
                            "type": "solid",
                            "opacity": 0.6,
                            "color": "source"
                        },
                        "tooltip": {"formatter": f"Shared: {', '.join(list(intersection)[:10])}"}
                    })
                    connection_counts[i] += 1
                    connection_counts[j] += 1
                    
                elif sem_score > 0.85: 
                    links.append({
                        "source": str(i), "target": str(j),
                        "value": round(sem_score, 2),
                        "lineStyle": {
                            "width": 1, "type": "dashed", "color": "#ccc", "opacity": 0.5
                        },
                        "tooltip": {"formatter": f"Semantic Similarity: {sem_score:.2f}"}
                    })

        for i in range(len(nodes)):
            nodes[i]["symbolSize"] = min(10 + (connection_counts[i] * 2), 70)

        return {"nodes": nodes, "links": links, "categories": categories}
    except Exception as e:
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"Graph Error: {str(e)}")